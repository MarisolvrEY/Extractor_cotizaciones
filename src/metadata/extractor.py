"""
meta.py  –  Extractor MÁXIMO de metadatos primarios
====================================================
Escanea una carpeta y vuelca TODOS los metadatos posibles de cada archivo
en un Excel con 3 hojas: Metadatos completos | Resumen | Leyenda.

Uso:
    python meta.py <carpeta_entrada> <carpeta_salida>
    python meta.py ./archivos ./resultado --recursivo
    python meta.py ./archivos ./resultado --recursivo --nombre mi_inventario

Tipos soportados:
    Imágenes  : jpg jpeg png tiff tif heic webp bmp gif raw cr2 nef arw dng
    PDF       : pdf
    Office    : docx xlsx pptx doc xls
    OpenDoc   : odt ods odp
    Audio     : mp3 flac ogg wav m4a aac wma opus aiff
    Video     : mp4 mov mkv avi webm flv wmv m4v 3gp
    Ebook     : epub
    Correo    : msg eml
    Web       : html htm
    Comprimido: zip
    Vectorial : svg
    Genérico  : cualquier otro → ExifTool si disponible

Librerías Python:
    pip install Pillow piexif pypdf python-docx openpyxl python-pptx mutagen xlrd ebooklib extract-msg

Programas externos (opcionales pero recomendados):
    FFmpeg  → videos          https://ffmpeg.org
    ExifTool → todo           https://exiftool.org
"""

import os, sys, json, argparse, subprocess, hashlib, mimetypes, struct, zipfile, re
from pathlib import Path
from datetime import datetime
from xml.etree import ElementTree as ET

# ── Imports opcionales ────────────────────────────────────────────────────────
def _try(module, pkg=None):
    try: return __import__(module)
    except ImportError:
        print(f"[AVISO] Falta librería → pip install {pkg or module}")
        return None

PIL      = _try("PIL",     "Pillow")
piexif   = _try("piexif")
pypdf    = _try("pypdf")
docx_m   = _try("docx",   "python-docx")
openpyxl = _try("openpyxl")
pptx_m   = _try("pptx",   "python-pptx")
mutagen  = _try("mutagen")
xlrd     = _try("xlrd")
ebooklib = _try("ebooklib")
extract_msg = _try("extract_msg")


# ═══════════════════════════════════════════════════════════════════════════════
# UTILIDADES
# ═══════════════════════════════════════════════════════════════════════════════

def _s(v):
    if isinstance(v, bytes):
        try: return v.decode("utf-8", errors="replace").strip("\x00")
        except: return v.hex()
    if isinstance(v, datetime): return v.isoformat()
    if isinstance(v, (list, tuple)): return " | ".join(_s(i) for i in v)
    if isinstance(v, dict): return json.dumps(v, ensure_ascii=False, default=str)
    return str(v).strip() if v is not None else ""

def _hash(path):
    md5 = hashlib.md5(); sha1 = hashlib.sha1(); sha256 = hashlib.sha256()
    try:
        with open(path, "rb") as f:
            for chunk in iter(lambda: f.read(65536), b""):
                md5.update(chunk); sha1.update(chunk); sha256.update(chunk)
        return md5.hexdigest(), sha1.hexdigest(), sha256.hexdigest()
    except: return "", "", ""

def _cmd_ok(cmd):
    try: subprocess.run([cmd, "-version"], capture_output=True, check=True); return True
    except: return False

def _ffprobe_ok():   return _cmd_ok("ffprobe")
def _exiftool_ok():
    try: subprocess.run(["exiftool","-ver"], capture_output=True, check=True); return True
    except: return False

def _gps_decimal(deg_list, ref):
    """Convierte coordenadas EXIF IFDRational a decimal."""
    try:
        def r(x):
            if hasattr(x, 'numerator'): return x.numerator / x.denominator
            if isinstance(x, tuple) and len(x)==2: return x[0]/x[1]
            return float(x)
        d = r(deg_list[0]) + r(deg_list[1])/60 + r(deg_list[2])/3600
        if ref in ("S","W"): d = -d
        return round(d, 8)
    except: return None


# ═══════════════════════════════════════════════════════════════════════════════
# 1. SISTEMA DE ARCHIVOS
# ═══════════════════════════════════════════════════════════════════════════════

def meta_filesystem(path):
    p  = Path(path)
    st = p.stat()
    md5, sha1, sha256 = _hash(path)
    m = {
        "fs_nombre_archivo"    : p.name,
        "fs_nombre_sin_ext"    : p.stem,
        "fs_extension"         : p.suffix.lower(),
        "fs_ruta_completa"     : str(p.resolve()),
        "fs_carpeta_padre"     : str(p.parent),
        "fs_tamaño_bytes"      : st.st_size,
        "fs_tamaño_kb"         : round(st.st_size / 1024, 2),
        "fs_tamaño_mb"         : round(st.st_size / 1_048_576, 4),
        "fs_fecha_modificacion": datetime.fromtimestamp(st.st_mtime).isoformat(),
        "fs_fecha_acceso"      : datetime.fromtimestamp(st.st_atime).isoformat(),
        "fs_tipo_mime"         : mimetypes.guess_type(path)[0] or "desconocido",
        "fs_md5"               : md5,
        "fs_sha1"              : sha1,
        "fs_sha256"            : sha256,
    }
    # ctime: en Windows = creación real; en Linux = cambio de inode
    try: m["fs_fecha_ctime"] = datetime.fromtimestamp(st.st_ctime).isoformat()
    except: pass
    return m


# ═══════════════════════════════════════════════════════════════════════════════
# 2. IMÁGENES (JPG, PNG, TIFF, HEIC, WEBP, BMP, GIF, RAW…)
# ═══════════════════════════════════════════════════════════════════════════════

def meta_imagen(path):
    m = {}
    if not PIL: return m
    from PIL import Image, ExifTags, TiffImagePlugin

    try:
        img = Image.open(path)
        m["img_formato"]       = img.format
        m["img_modo_color"]    = img.mode
        m["img_ancho_px"]      = img.width
        m["img_alto_px"]       = img.height
        m["img_megapixeles"]   = round(img.width * img.height / 1_000_000, 2)
        m["img_relacion_aspecto"] = f"{img.width}:{img.height}"

        info = img.info or {}

        # DPI / resolución
        dpi = info.get("dpi") or info.get("resolution")
        if dpi:
            try: m["img_dpi"] = f"{round(float(dpi[0]),1)} x {round(float(dpi[1]),1)}"
            except: m["img_dpi"] = _s(dpi)

        # Profundidad de bits
        for k in ["bits","n_frames","loop","version","duration"]:
            if info.get(k) is not None: m[f"img_{k}"] = _s(info[k])

        # Comentario embebido
        for k in ["comment","Comment","artist","copyright"]:
            if info.get(k): m[f"img_info_{k}"] = _s(info[k])

        # Perfil ICC → nombre del espacio de color
        if "icc_profile" in info:
            icc = info["icc_profile"]
            m["img_icc_tamaño_bytes"] = len(icc)
            try:
                # Los primeros 128 bytes del ICC contienen info legible
                m["img_icc_espacio_color"] = icc[16:20].decode("ascii","replace").strip()
                m["img_icc_clase"]         = icc[12:16].decode("ascii","replace").strip()
                desc_offset = struct.unpack(">I", icc[128:132])[0]
                m["img_icc_descripcion"]   = icc[128:128+60].decode("latin-1","replace").strip()
            except: pass

        # PNG chunks adicionales (tEXt / iTXt / zTXt / tIME)
        if img.format == "PNG":
            for key, val in info.items():
                if key.lower() not in ("dpi","icc_profile","bits") and isinstance(val, str):
                    m[f"img_png_{key}"] = val
            # tIME chunk → fecha de última modificación embebida en el PNG
            if hasattr(img, "_getxmp"):
                try:
                    xmp = img._getxmp()
                    if xmp: m["img_xmp_raw"] = str(xmp)[:500]
                except: pass

        # GIF info
        if img.format == "GIF":
            m["img_gif_frames"]  = getattr(img, "n_frames", 1)
            m["img_gif_loop"]    = info.get("loop", 0)
            m["img_gif_version"] = info.get("version","")

        # EXIF completo vía PIL
        raw_exif = None
        if hasattr(img, "_getexif"):
            try: raw_exif = img._getexif()
            except: pass

        if raw_exif:
            for tid, val in raw_exif.items():
                name = ExifTags.TAGS.get(tid, f"tag_{tid}")
                if name in ("MakerNote","UserComment") and isinstance(val, bytes):
                    m[f"exif_{name}_hex"] = val[:64].hex()
                    continue
                if name == "GPSInfo" and isinstance(val, dict):
                    # Parsear GPS completo
                    gps = val
                    GPS = ExifTags.GPSTAGS
                    gps_named = {GPS.get(k, str(k)): v for k, v in gps.items()}
                    for gk, gv in gps_named.items():
                        m[f"exif_GPS_{gk}"] = _s(gv)
                    # Convertir a decimal
                    try:
                        lat = _gps_decimal(gps_named.get("GPSLatitude",[]),
                                           gps_named.get("GPSLatitudeRef","N"))
                        lon = _gps_decimal(gps_named.get("GPSLongitude",[]),
                                           gps_named.get("GPSLongitudeRef","E"))
                        if lat and lon:
                            m["exif_GPS_latitud_decimal"]  = lat
                            m["exif_GPS_longitud_decimal"] = lon
                            m["exif_GPS_google_maps"] = f"https://maps.google.com/?q={lat},{lon}"
                    except: pass
                    continue
                m[f"exif_{name}"] = _s(val)

        # EXIF vía piexif (más bajo nivel, IFDs completos)
        elif "exif" in info and piexif:
            try:
                exif = piexif.load(info["exif"])
                for ifd, tags in exif.items():
                    if not isinstance(tags, dict): continue
                    for tid, val in tags.items():
                        try: name = piexif.TAGS[ifd][tid]["name"]
                        except: name = str(tid)
                        if name == "MakerNote": continue
                        m[f"exif_{ifd}_{name}"] = _s(val)
            except: pass

        # IPTC (si disponible)
        try:
            from PIL import IptcImagePlugin
            iptc = IptcImagePlugin.getiptcinfo(img) or {}
            IPTC_TAGS = {
                (2,5): "titulo_objeto", (2,10): "urgencia", (2,15): "categoria",
                (2,20): "suplemento_categoria", (2,25): "palabras_clave",
                (2,40): "instrucciones_especiales", (2,55): "fecha_creacion",
                (2,60): "hora_creacion", (2,65): "tiempo_digital",
                (2,80): "autor", (2,85): "titulo_autor",
                (2,90): "ciudad", (2,92): "sublocalidad",
                (2,95): "estado", (2,100): "codigo_pais",
                (2,101): "pais", (2,103): "referencia_transmision",
                (2,105): "encabezado", (2,110): "credito",
                (2,115): "fuente", (2,116): "copyright",
                (2,120): "descripcion", (2,122): "autor_descripcion",
            }
            for k, v in iptc.items():
                label = IPTC_TAGS.get(k, f"iptc_{k[0]}_{k[1]}")
                m[f"iptc_{label}"] = _s(v)
        except: pass

    except Exception as e:
        m["img_error"] = str(e)

    return m


# ═══════════════════════════════════════════════════════════════════════════════
# 3. PDF
# ═══════════════════════════════════════════════════════════════════════════════

def meta_pdf(path):
    m = {}
    if not pypdf: return m
    try:
        import pypdf as pp
        r = pp.PdfReader(path)

        # Info básica
        info = r.metadata or {}
        for attr in ["title","author","subject","creator","producer"]:
            v = getattr(info, attr, None)
            if v: m[f"pdf_{attr}"] = _s(v)

        m["pdf_fecha_creacion"]        = _s(getattr(info, "creation_date", None))
        m["pdf_fecha_modificacion"]    = _s(getattr(info, "modification_date", None))
        m["pdf_num_paginas"]           = len(r.pages)
        m["pdf_encriptado"]            = r.is_encrypted

        # Versión PDF
        try: m["pdf_version"] = r.pdf_header
        except: pass

        # Todos los campos raw del diccionario /Info
        for k, v in info.items():
            clean_k = str(k).lstrip("/").lower()
            if clean_k not in ("title","author","subject","creator","producer",
                               "creationdate","moddate"):
                m[f"pdf_info_{clean_k}"] = _s(v)

        # Tamaño de página (primera página)
        try:
            page = r.pages[0]
            box  = page.mediabox
            m["pdf_pag_ancho_pt"]  = round(float(box.width), 2)
            m["pdf_pag_alto_pt"]   = round(float(box.height), 2)
            m["pdf_pag_ancho_mm"]  = round(float(box.width) * 0.352778, 1)
            m["pdf_pag_alto_mm"]   = round(float(box.height) * 0.352778, 1)
        except: pass

        # Permisos
        try:
            perms = r._reader.get_fields()
            if perms: m["pdf_tiene_formulario"] = True
        except: pass

        # XMP completo
        xmp = r.xmp_metadata
        if xmp:
            xmp_fields = {
                "dc_title": "xmp_dc_titulo",
                "dc_creator": "xmp_dc_creador",
                "dc_description": "xmp_dc_descripcion",
                "dc_subject": "xmp_dc_materia",
                "dc_rights": "xmp_dc_derechos",
                "dc_format": "xmp_dc_formato",
                "dc_identifier": "xmp_dc_identificador",
                "dc_language": "xmp_dc_idioma",
                "dc_publisher": "xmp_dc_editorial",
                "dc_relation": "xmp_dc_relacion",
                "dc_source": "xmp_dc_fuente",
                "dc_type": "xmp_dc_tipo",
                "xmp_create_date": "xmp_fecha_creacion",
                "xmp_modify_date": "xmp_fecha_modificacion",
                "xmp_metadata_date": "xmp_fecha_metadata",
                "xmp_creator_tool": "xmp_herramienta_creacion",
                "xmp_label": "xmp_etiqueta",
                "xmp_rating": "xmp_calificacion",
            }
            for attr, label in xmp_fields.items():
                try:
                    v = getattr(xmp, attr, None)
                    if v: m[label] = _s(v)
                except: pass

            # PDF/A info
            try:
                raw = xmp.getElement("", "pdfaid", "part")
                if raw: m["xmp_pdfa_parte"] = _s(raw)
            except: pass

        # Fuentes embebidas (primeras 20)
        try:
            fonts = set()
            for page in r.pages[:20]:
                res = page.get("/Resources",{})
                font_dict = res.get("/Font",{})
                for fname, fobj in font_dict.items():
                    try:
                        obj = fobj.get_object()
                        base = obj.get("/BaseFont","")
                        if base: fonts.add(str(base).lstrip("/"))
                    except: pass
            if fonts:
                m["pdf_fuentes_embebidas"] = " | ".join(sorted(fonts)[:30])
                m["pdf_num_fuentes"] = len(fonts)
        except: pass

        # Adjuntos
        try:
            catalog = r.trailer["/Root"]
            names   = catalog.get("/Names",{})
            embeds  = names.get("/EmbeddedFiles",{})
            if embeds:
                m["pdf_tiene_adjuntos"] = True
        except: pass

    except Exception as e:
        m["pdf_error"] = str(e)
    return m


# ═══════════════════════════════════════════════════════════════════════════════
# 4. OFFICE (DOCX / XLSX / PPTX) — core.xml + app.xml
# ═══════════════════════════════════════════════════════════════════════════════

def _parse_office_zip(path, prefix):
    """Extrae core.xml y app.xml de cualquier archivo Office (son ZIPs)."""
    m = {}
    try:
        with zipfile.ZipFile(path, "r") as z:
            names = z.namelist()

            # ── core.xml ────────────────────────────────────────────────────
            core_candidates = [n for n in names if "core.xml" in n.lower()]
            if core_candidates:
                xml = z.read(core_candidates[0])
                root = ET.fromstring(xml)
                NS = {
                    "cp":"http://schemas.openxmlformats.org/package/2006/metadata/core-properties",
                    "dc":"http://purl.org/dc/elements/1.1/",
                    "dcterms":"http://purl.org/dc/terms/",
                    "xsi":"http://www.w3.org/2001/XMLSchema-instance",
                }
                field_map = {
                    "dc:creator":             f"{prefix}_autor_creador",
                    "cp:lastModifiedBy":       f"{prefix}_ultimo_editor",
                    "dc:title":               f"{prefix}_titulo",
                    "dc:subject":             f"{prefix}_asunto",
                    "dc:description":         f"{prefix}_descripcion",
                    "cp:keywords":            f"{prefix}_palabras_clave",
                    "cp:category":            f"{prefix}_categoria",
                    "cp:contentStatus":       f"{prefix}_estado",
                    "dc:language":            f"{prefix}_idioma",
                    "cp:revision":            f"{prefix}_revision",
                    "cp:version":             f"{prefix}_version",
                    "cp:identifier":          f"{prefix}_identificador",
                    "dcterms:created":        f"{prefix}_fecha_creacion",
                    "dcterms:modified":       f"{prefix}_fecha_modificacion",
                    "cp:lastPrinted":         f"{prefix}_ultima_impresion",
                    "cp:contentType":         f"{prefix}_tipo_contenido",
                }
                for tag, label in field_map.items():
                    ns, local = tag.split(":") if ":" in tag else ("", tag)
                    uri = NS.get(ns, "")
                    el = root.find(f"{{{uri}}}{local}" if uri else local)
                    if el is not None and el.text:
                        m[label] = el.text.strip()

            # ── app.xml ─────────────────────────────────────────────────────
            app_candidates = [n for n in names if "app.xml" in n.lower()]
            if app_candidates:
                xml = z.read(app_candidates[0])
                root = ET.fromstring(xml)
                NS_APP = "http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"
                NS_VT  = "http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes"
                app_fields = [
                    "Application","AppVersion","Company","Manager","Template",
                    "Words","Characters","CharactersWithSpaces","Lines","Paragraphs",
                    "Pages","Slides","Notes","HiddenSlides","MMClips",
                    "TotalTime","LinksUpToDate","SharedDoc","HyperlinksChanged",
                    "ScaleCrop","HeadingPairs","TitlesOfParts",
                ]
                for field in app_fields:
                    el = root.find(f"{{{NS_APP}}}{field}")
                    if el is not None and el.text:
                        m[f"{prefix}_app_{field.lower()}"] = el.text.strip()

            # ── Archivos internos (lista de relaciones) ──────────────────────
            m[f"{prefix}_archivos_internos"] = str(len(names))

    except Exception as e:
        m[f"{prefix}_zip_error"] = str(e)
    return m


def meta_docx(path):
    m = _parse_office_zip(path, "docx")
    if not docx_m: return m
    try:
        from docx import Document
        doc = Document(path)
        p = doc.core_properties
        # Core props vía python-docx (redundante pero por si acaso)
        for attr in ["title","author","subject","keywords","description","category",
                     "content_status","identifier","language","last_modified_by",
                     "last_printed","created","modified","revision","version","comments"]:
            v = getattr(p, attr, None)
            if v and f"docx_{attr}" not in m:
                m[f"docx_{attr}"] = _s(v)

        m["docx_num_parrafos"]  = len(doc.paragraphs)
        m["docx_num_tablas"]    = len(doc.tables)
        m["docx_num_secciones"] = len(doc.sections)
        m["docx_num_imagenes"]  = sum(1 for _ in doc.inline_shapes)

        # Texto total (primeras 500 chars)
        texto = " ".join(p.text for p in doc.paragraphs if p.text.strip())
        m["docx_preview_texto"] = texto[:500]
        m["docx_num_caracteres_aprox"] = len(texto)

        # Estilos usados
        estilos = list({p.style.name for p in doc.paragraphs if p.style})
        m["docx_estilos_usados"] = " | ".join(sorted(estilos)[:15])

        # Comentarios
        try:
            from docx.oxml.ns import qn
            comentarios = doc.element.body.findall(
                f".//{{{doc.element.body.nsmap.get('w','')}}}"+"comment", {})
            m["docx_num_comentarios"] = len(comentarios)
        except: pass

    except Exception as e:
        m["docx_error"] = str(e)
    return m


def meta_xlsx(path):
    m = _parse_office_zip(path, "xlsx")
    if not openpyxl: return m
    try:
        import openpyxl as ox
        wb = ox.load_workbook(path, read_only=True, data_only=True)
        p  = wb.properties
        for attr in ["title","subject","creator","keywords","description","lastModifiedBy",
                     "revision","created","modified","category","contentStatus",
                     "language","version","identifier"]:
            v = getattr(p, attr, None)
            if v and f"xlsx_{attr}" not in m:
                m[f"xlsx_{attr}"] = _s(v)

        m["xlsx_hojas"]         = " | ".join(wb.sheetnames)
        m["xlsx_num_hojas"]     = len(wb.sheetnames)

        # Estadísticas por hoja
        total_celdas = 0
        for sname in wb.sheetnames[:10]:
            ws = wb[sname]
            try:
                filas = ws.max_row or 0
                cols  = ws.max_column or 0
                m[f"xlsx_hoja_{sname[:20]}_filas"] = filas
                m[f"xlsx_hoja_{sname[:20]}_cols"]  = cols
                total_celdas += (filas or 0) * (cols or 0)
            except: pass
        m["xlsx_total_celdas_aprox"] = total_celdas
        wb.close()
    except Exception as e:
        m["xlsx_error"] = str(e)
    return m


def meta_pptx(path):
    m = _parse_office_zip(path, "pptx")
    if not pptx_m: return m
    try:
        from pptx import Presentation
        prs = Presentation(path)
        p = prs.core_properties
        for attr in ["title","author","subject","keywords","description","category",
                     "content_status","identifier","language","last_modified_by",
                     "last_printed","created","modified","revision","version","comments"]:
            v = getattr(p, attr, None)
            if v and f"pptx_{attr}" not in m:
                m[f"pptx_{attr}"] = _s(v)

        m["pptx_num_diapositivas"]  = len(prs.slides)
        m["pptx_ancho_pt"]  = round(prs.slide_width.pt, 1)  if prs.slide_width  else ""
        m["pptx_alto_pt"]   = round(prs.slide_height.pt, 1) if prs.slide_height else ""

        # Texto de la primera diapositiva
        try:
            shapes = prs.slides[0].shapes
            texts  = [s.text for s in shapes if hasattr(s,"text") and s.text.strip()]
            m["pptx_texto_diap1"] = " | ".join(texts)[:300]
        except: pass

        # Layouts disponibles
        try: m["pptx_layouts"] = str(len(prs.slide_layouts))
        except: pass

    except Exception as e:
        m["pptx_error"] = str(e)
    return m


# ── DOC / XLS (formatos binarios antiguos) ───────────────────────────────────
def meta_doc_old(path):
    """DOC binario — metadata vía ExifTool principalmente."""
    m = {}
    # Intentar leer summary information stream (OLE2)
    try:
        import struct
        with open(path,"rb") as f:
            sig = f.read(8)
        if sig == b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1":
            m["doc_formato"] = "OLE2 / Compound Document (Word 97-2003)"
        else:
            m["doc_formato"] = "Desconocido"
    except: pass
    return m

def meta_xls_old(path):
    """XLS binario — via xlrd."""
    m = {}
    if not xlrd: return m
    try:
        wb = xlrd.open_workbook(path)
        m["xls_hojas"]     = " | ".join(wb.sheet_names())
        m["xls_num_hojas"] = wb.nsheets
        m["xls_autor"]     = wb.user_name or ""
        # Propiedades de documento si están disponibles
        try:
            for attr in ["user_name","datemode"]:
                v = getattr(wb, attr, None)
                if v is not None: m[f"xls_{attr}"] = _s(v)
        except: pass
    except Exception as e:
        m["xls_error"] = str(e)
    return m


# ═══════════════════════════════════════════════════════════════════════════════
# 5. OPENDOCUMENT (ODT / ODS / ODP)
# ═══════════════════════════════════════════════════════════════════════════════

def meta_opendoc(path, prefix):
    """ODT / ODS / ODP son ZIPs con meta.xml y manifest.xml."""
    m = {}
    try:
        with zipfile.ZipFile(path,"r") as z:
            names = z.namelist()
            m[f"{prefix}_archivos_internos"] = str(len(names))

            # meta.xml
            if "meta.xml" in names:
                xml = z.read("meta.xml")
                root = ET.fromstring(xml)
                OFFICE = "urn:oasis:names:tc:opendocument:xmlns:office:1.0"
                META   = "urn:oasis:names:tc:opendocument:xmlns:meta:1.0"
                DC     = "http://purl.org/dc/elements/1.1/"

                stats = root.find(f"{{{OFFICE}}}meta/{{{META}}}document-statistic")
                if stats is not None:
                    for attr_name, attr_val in stats.attrib.items():
                        local = attr_name.split("}")[-1] if "}" in attr_name else attr_name
                        m[f"{prefix}_stat_{local}"] = attr_val

                field_map = {
                    f"{{{DC}}}title":          f"{prefix}_titulo",
                    f"{{{DC}}}description":    f"{prefix}_descripcion",
                    f"{{{DC}}}subject":        f"{prefix}_asunto",
                    f"{{{DC}}}creator":        f"{prefix}_creador",
                    f"{{{DC}}}language":       f"{prefix}_idioma",
                    f"{{{META}}}creation-date":f"{prefix}_fecha_creacion",
                    f"{{{META}}}date":         f"{prefix}_fecha_modificacion",
                    f"{{{META}}}editing-duration": f"{prefix}_tiempo_edicion",
                    f"{{{META}}}editing-cycles":   f"{prefix}_ciclos_edicion",
                    f"{{{META}}}generator":    f"{prefix}_aplicacion_creadora",
                    f"{{{META}}}keyword":      f"{prefix}_palabras_clave",
                    f"{{{META}}}initial-creator": f"{prefix}_creador_original",
                    f"{{{META}}}printed-by":      f"{prefix}_impreso_por",
                    f"{{{META}}}print-date":       f"{prefix}_fecha_impresion",
                }
                meta_node = root.find(f"{{{OFFICE}}}meta")
                if meta_node is not None:
                    for child in meta_node:
                        label = field_map.get(child.tag)
                        if label and child.text:
                            m[label] = child.text.strip()
    except Exception as e:
        m[f"{prefix}_error"] = str(e)
    return m


# ═══════════════════════════════════════════════════════════════════════════════
# 6. AUDIO
# ═══════════════════════════════════════════════════════════════════════════════

def meta_audio(path):
    m = {}
    if not mutagen: return m
    try:
        import mutagen as mg

        # ── Info técnica (modo raw, no easy) ─────────────────────────────────
        audio_raw = mg.File(path)
        if audio_raw is not None:
            info = audio_raw.info
            m["audio_duracion_seg"]     = round(info.length, 3) if hasattr(info,"length") else ""
            m["audio_duracion_hms"]     = str(datetime.utcfromtimestamp(
                                            info.length).strftime("%H:%M:%S")) if hasattr(info,"length") else ""
            m["audio_bitrate_bps"]      = getattr(info, "bitrate", "")
            m["audio_bitrate_kbps"]     = round(getattr(info,"bitrate",0)/1000, 1) if getattr(info,"bitrate",0) else ""
            m["audio_canales"]          = getattr(info, "channels", "")
            m["audio_sample_rate_hz"]   = getattr(info, "sample_rate", "")
            m["audio_bits_por_muestra"] = getattr(info, "bits_per_sample", "")
            m["audio_encoder_info"]     = getattr(info, "encoder_info", "")
            m["audio_modo_stereo"]      = getattr(info, "mode", "")  # 0=stereo, 1=joint stereo, etc.
            m["audio_codec"]            = type(info).__name__

            # Tags raw — TODOS los frames/campos disponibles
            for key, val in audio_raw.items():
                clean_key = re.sub(r"[^a-zA-Z0-9_]", "_", str(key)).lower()
                # Saltar portada binaria
                if "apic" in clean_key or "covr" in clean_key or "cover" in clean_key:
                    m["audio_tiene_portada_embebida"] = True
                    continue
                m[f"audio_tag_{clean_key}"] = _s(val)

        # ── Easy tags (más limpios para campos estándar) ──────────────────────
        audio_easy = mg.File(path, easy=True)
        if audio_easy:
            easy_fields = [
                "title","artist","album","date","genre","tracknumber","albumartist",
                "composer","lyricist","conductor","remixer","performer","arranger",
                "author","encoder","discnumber","comment","lyrics","language",
                "copyright","license","isrc","barcode","catalognumber","label",
                "bpm","key","media","organization","version","website",
                "musicbrainz_trackid","musicbrainz_albumid","musicbrainz_artistid",
                "musicbrainz_albumartistid","musicbrainz_releasegroupid",
                "replaygain_track_gain","replaygain_track_peak",
                "replaygain_album_gain","replaygain_album_peak",
            ]
            for tag in easy_fields:
                v = audio_easy.get(tag)
                if v:
                    m[f"audio_{tag}"] = _s(v[0] if isinstance(v,list) else v)

    except Exception as e:
        m["audio_error"] = str(e)
    return m


# ═══════════════════════════════════════════════════════════════════════════════
# 7. VIDEO (ffprobe)
# ═══════════════════════════════════════════════════════════════════════════════

def meta_video(path):
    m = {}
    if not _ffprobe_ok():
        m["video_error"] = "ffprobe no disponible — instala FFmpeg: https://ffmpeg.org"
        return m
    try:
        # Formato + streams + capítulos + frames de referencia
        r = subprocess.run(
            ["ffprobe","-v","quiet","-print_format","json",
             "-show_format","-show_streams","-show_chapters", str(path)],
            capture_output=True, text=True, check=True
        )
        info = json.loads(r.stdout)
        fmt  = info.get("format", {})
        tags = fmt.get("tags", {})

        m["video_formato_corto"]     = fmt.get("format_name","")
        m["video_formato_largo"]     = fmt.get("format_long_name","")
        m["video_duracion_seg"]      = round(float(fmt["duration"]),3) if "duration" in fmt else ""
        m["video_duracion_hms"]      = str(datetime.utcfromtimestamp(
                                            float(fmt.get("duration",0))).strftime("%H:%M:%S"))
        m["video_tamaño_bytes"]      = fmt.get("size","")
        m["video_bitrate_total_bps"] = fmt.get("bit_rate","")
        m["video_num_streams"]       = fmt.get("nb_streams","")
        m["video_num_programas"]     = fmt.get("nb_programs","")

        # Todos los tags del contenedor (fecha creación, título, encoder, GPS, etc.)
        for k, v in tags.items():
            m[f"video_tag_{k.lower()}"] = _s(v)

        # Streams detallados
        video_idx = audio_idx = sub_idx = 0
        for stream in info.get("streams", []):
            stype = stream.get("codec_type","?")
            idx   = stream.get("index","?")

            if stype == "video":
                pfx = f"video_v{video_idx}"
                video_idx += 1
                m[f"{pfx}_codec"]           = stream.get("codec_name","")
                m[f"{pfx}_codec_largo"]      = stream.get("codec_long_name","")
                m[f"{pfx}_perfil"]          = stream.get("profile","")
                m[f"{pfx}_nivel"]           = stream.get("level","")
                m[f"{pfx}_ancho_px"]        = stream.get("width","")
                m[f"{pfx}_alto_px"]         = stream.get("height","")
                m[f"{pfx}_ancho_codificado"]= stream.get("coded_width","")
                m[f"{pfx}_alto_codificado"] = stream.get("coded_height","")
                m[f"{pfx}_sar"]             = stream.get("sample_aspect_ratio","")
                m[f"{pfx}_dar"]             = stream.get("display_aspect_ratio","")
                m[f"{pfx}_pix_fmt"]         = stream.get("pix_fmt","")
                m[f"{pfx}_fps"]             = stream.get("r_frame_rate","")
                m[f"{pfx}_fps_avg"]         = stream.get("avg_frame_rate","")
                m[f"{pfx}_color_range"]     = stream.get("color_range","")
                m[f"{pfx}_color_space"]     = stream.get("color_space","")
                m[f"{pfx}_color_transfer"]  = stream.get("color_transfer","")
                m[f"{pfx}_color_primaries"] = stream.get("color_primaries","")
                m[f"{pfx}_bitrate"]         = stream.get("bit_rate","")
                m[f"{pfx}_bits_raw"]        = stream.get("bits_per_raw_sample","")
                m[f"{pfx}_num_frames"]      = stream.get("nb_frames","")
                # Rotación
                stags = stream.get("tags",{})
                if "rotate" in stags: m[f"{pfx}_rotacion"] = stags["rotate"]
                m[f"{pfx}_idioma"] = stags.get("language","")

            elif stype == "audio":
                pfx = f"video_a{audio_idx}"
                audio_idx += 1
                m[f"{pfx}_codec"]           = stream.get("codec_name","")
                m[f"{pfx}_codec_largo"]      = stream.get("codec_long_name","")
                m[f"{pfx}_perfil"]          = stream.get("profile","")
                m[f"{pfx}_sample_rate"]     = stream.get("sample_rate","")
                m[f"{pfx}_canales"]         = stream.get("channels","")
                m[f"{pfx}_layout"]          = stream.get("channel_layout","")
                m[f"{pfx}_bits"]            = stream.get("bits_per_sample","")
                m[f"{pfx}_bitrate"]         = stream.get("bit_rate","")
                stags = stream.get("tags",{})
                m[f"{pfx}_idioma"]          = stags.get("language","")
                m[f"{pfx}_titulo"]          = stags.get("title","")

            elif stype == "subtitle":
                pfx = f"video_s{sub_idx}"
                sub_idx += 1
                m[f"{pfx}_codec"]   = stream.get("codec_name","")
                stags = stream.get("tags",{})
                m[f"{pfx}_idioma"]  = stags.get("language","")
                m[f"{pfx}_titulo"]  = stags.get("title","")

        m["video_num_pistas_video"]    = video_idx
        m["video_num_pistas_audio"]    = audio_idx
        m["video_num_pistas_subtitulo"]= sub_idx

        # Capítulos
        chapters = info.get("chapters", [])
        if chapters:
            m["video_num_capitulos"] = len(chapters)
            caps = []
            for ch in chapters[:10]:
                t = ch.get("tags",{}).get("title","")
                s = ch.get("start_time","")
                caps.append(f"{t}@{s}s" if t else f"cap@{s}s")
            m["video_capitulos"] = " | ".join(caps)

    except Exception as e:
        m["video_error"] = str(e)
    return m


# ═══════════════════════════════════════════════════════════════════════════════
# 8. EPUB
# ═══════════════════════════════════════════════════════════════════════════════

def meta_epub(path):
    m = {}
    try:
        with zipfile.ZipFile(path,"r") as z:
            names = z.namelist()
            m["epub_archivos_internos"] = len(names)

            # Buscar OPF (contenedor principal de metadata)
            opf_candidates = [n for n in names if n.endswith(".opf")]
            if not opf_candidates:
                # Leer container.xml
                if "META-INF/container.xml" in names:
                    cont_xml = z.read("META-INF/container.xml")
                    root = ET.fromstring(cont_xml)
                    for rf in root.iter("{urn:oasis:names:tc:opendocument:xmlns:container}rootfile"):
                        opf_candidates.append(rf.get("full-path",""))

            if opf_candidates:
                opf_xml = z.read(opf_candidates[0])
                root = ET.fromstring(opf_xml)
                DC = "http://purl.org/dc/elements/1.1/"
                OPF = "http://www.idpf.org/2007/opf"

                field_map = {
                    f"{{{DC}}}title":       "epub_titulo",
                    f"{{{DC}}}creator":     "epub_autor",
                    f"{{{DC}}}publisher":   "epub_editorial",
                    f"{{{DC}}}date":        "epub_fecha",
                    f"{{{DC}}}language":    "epub_idioma",
                    f"{{{DC}}}description": "epub_descripcion",
                    f"{{{DC}}}subject":     "epub_materia",
                    f"{{{DC}}}rights":      "epub_derechos",
                    f"{{{DC}}}identifier":  "epub_identificador",
                    f"{{{DC}}}source":      "epub_fuente",
                    f"{{{DC}}}contributor": "epub_contribuidor",
                    f"{{{DC}}}type":        "epub_tipo",
                    f"{{{DC}}}format":      "epub_formato",
                    f"{{{DC}}}coverage":    "epub_cobertura",
                    f"{{{DC}}}relation":    "epub_relacion",
                }
                metadata_node = root.find(f"{{{OPF}}}metadata")
                if metadata_node is not None:
                    for child in metadata_node:
                        label = field_map.get(child.tag)
                        if label and child.text:
                            existing = m.get(label)
                            m[label] = f"{existing} | {child.text.strip()}" if existing else child.text.strip()

                # Meta tags adicionales (modified, cover-image, etc.)
                if metadata_node is not None:
                    for meta in metadata_node.findall(f"{{{OPF}}}meta"):
                        name = meta.get("name") or meta.get("property","")
                        val  = meta.get("content") or meta.text or ""
                        if name and val:
                            clean = re.sub(r"[^a-z0-9_]","_",name.lower())
                            m[f"epub_meta_{clean}"] = val.strip()

                # Spine (orden de lectura)
                spine = root.find(f"{{{OPF}}}spine")
                if spine is not None:
                    m["epub_num_capitulos"] = len(list(spine))

    except Exception as e:
        m["epub_error"] = str(e)
    return m


# ═══════════════════════════════════════════════════════════════════════════════
# 9. CORREO ELECTRÓNICO (MSG / EML)
# ═══════════════════════════════════════════════════════════════════════════════

def meta_msg(path):
    m = {}
    if not extract_msg: return {"msg_error": "pip install extract-msg"}
    try:
        import extract_msg as em
        msg = em.Message(path)
        field_map = {
            "subject":      "msg_asunto",
            "sender":       "msg_remitente",
            "to":           "msg_destinatarios",
            "cc":           "msg_cc",
            "bcc":          "msg_bcc",
            "date":         "msg_fecha_envio",
            "messageId":    "msg_message_id",
            "importance":   "msg_importancia",
            "sensitivity":  "msg_sensibilidad",
            "categories":   "msg_categorias",
            "body":         None,  # omitir cuerpo completo
        }
        for attr, label in field_map.items():
            if label is None: continue
            v = getattr(msg, attr, None)
            if v: m[label] = _s(v)[:500]

        m["msg_num_adjuntos"] = len(msg.attachments)
        if msg.attachments:
            m["msg_adjuntos"] = " | ".join(
                a.longFilename or a.shortFilename or "adjunto"
                for a in msg.attachments[:10]
            )
        m["msg_preview_cuerpo"] = (_s(msg.body) or "")[:300]

    except Exception as e:
        m["msg_error"] = str(e)
    return m


def meta_eml(path):
    m = {}
    try:
        import email
        from email import policy as ep
        with open(path,"rb") as f:
            msg = email.message_from_binary_file(f, policy=ep.default)

        headers = [
            ("Subject",      "eml_asunto"),
            ("From",         "eml_remitente"),
            ("To",           "eml_destinatarios"),
            ("Cc",           "eml_cc"),
            ("Bcc",          "eml_bcc"),
            ("Date",         "eml_fecha_envio"),
            ("Message-ID",   "eml_message_id"),
            ("Reply-To",     "eml_reply_to"),
            ("X-Mailer",     "eml_cliente_correo"),
            ("X-Originating-IP","eml_ip_origen"),
            ("Content-Type", "eml_content_type"),
            ("MIME-Version", "eml_mime_version"),
            ("Importance",   "eml_importancia"),
            ("X-Priority",   "eml_prioridad"),
            ("Return-Path",  "eml_return_path"),
        ]
        for header, label in headers:
            v = msg.get(header)
            if v: m[label] = _s(v)[:500]

        # Adjuntos
        adjuntos = []
        for part in msg.walk():
            cd = part.get_content_disposition()
            if cd == "attachment":
                adjuntos.append(part.get_filename() or "adjunto")
        m["eml_num_adjuntos"] = len(adjuntos)
        if adjuntos: m["eml_adjuntos"] = " | ".join(adjuntos[:10])

        # Preview del body
        try:
            body = msg.get_body(preferencelist=("plain",))
            if body:
                m["eml_preview_cuerpo"] = body.get_content()[:300]
        except: pass

        # Recibido (ruta de servidores)
        received = msg.get_all("Received") or []
        m["eml_num_saltos_servidor"] = len(received)
        if received: m["eml_primer_servidor"] = received[-1][:200]

    except Exception as e:
        m["eml_error"] = str(e)
    return m


# ═══════════════════════════════════════════════════════════════════════════════
# 10. HTML / HTM
# ═══════════════════════════════════════════════════════════════════════════════

def meta_html(path):
    m = {}
    try:
        with open(path,"r", encoding="utf-8", errors="replace") as f:
            content = f.read()

        # Meta tags
        meta_re = re.finditer(
            r'<meta\s+([^>]+)>', content, re.IGNORECASE)
        for match in meta_re:
            attrs = match.group(1)
            name  = re.search(r'name=["\']([^"\']+)["\']', attrs, re.I)
            prop  = re.search(r'property=["\']([^"\']+)["\']', attrs, re.I)
            cont  = re.search(r'content=["\']([^"\']+)["\']', attrs, re.I)
            http_eq = re.search(r'http-equiv=["\']([^"\']+)["\']', attrs, re.I)
            key   = (name or prop or http_eq)
            if key and cont:
                clean = re.sub(r"[^a-z0-9_]","_", key.group(1).lower())
                m[f"html_meta_{clean}"] = cont.group(1)[:300]

        # Title
        title = re.search(r'<title[^>]*>([^<]+)</title>', content, re.I)
        if title: m["html_titulo"] = title.group(1).strip()

        # Charset
        charset = re.search(r'charset=["\']?([a-z0-9\-]+)', content, re.I)
        if charset: m["html_charset"] = charset.group(1)

        # Estadísticas
        m["html_longitud_chars"] = len(content)
        m["html_num_enlaces"]    = len(re.findall(r'<a\s', content, re.I))
        m["html_num_imagenes"]   = len(re.findall(r'<img\s', content, re.I))
        m["html_num_scripts"]    = len(re.findall(r'<script', content, re.I))
        m["html_doctype"]        = "HTML5" if "<!DOCTYPE html>" in content[:200].upper() else "Otro"

    except Exception as e:
        m["html_error"] = str(e)
    return m


# ═══════════════════════════════════════════════════════════════════════════════
# 11. ZIP
# ═══════════════════════════════════════════════════════════════════════════════

def meta_zip(path):
    m = {}
    try:
        with zipfile.ZipFile(path,"r") as z:
            infos = z.infolist()
            m["zip_num_archivos"]       = len(infos)
            m["zip_comentario"]         = _s(z.comment)
            m["zip_tamaño_total_sin_comprimir"] = sum(i.file_size for i in infos)
            m["zip_tamaño_total_comprimido"]    = sum(i.compress_size for i in infos)
            if m["zip_tamaño_total_sin_comprimir"]:
                ratio = 1 - m["zip_tamaño_total_comprimido"] / m["zip_tamaño_total_sin_comprimir"]
                m["zip_ratio_compresion_pct"] = round(ratio * 100, 1)

            # Archivos más pesados (top 5)
            top = sorted(infos, key=lambda i: i.file_size, reverse=True)[:5]
            m["zip_archivos_mas_grandes"] = " | ".join(
                f"{i.filename}({round(i.file_size/1024,1)}KB)" for i in top)

            # Fecha más reciente embebida
            dates = []
            for i in infos:
                try:
                    d = datetime(*i.date_time)
                    dates.append(d)
                except: pass
            if dates:
                m["zip_fecha_mas_reciente"] = max(dates).isoformat()
                m["zip_fecha_mas_antigua"]  = min(dates).isoformat()

            # Tipos de archivo contenidos
            exts = set(Path(i.filename).suffix.lower() for i in infos if "." in i.filename)
            m["zip_extensiones_contenidas"] = " ".join(sorted(exts)[:20])

    except Exception as e:
        m["zip_error"] = str(e)
    return m


# ═══════════════════════════════════════════════════════════════════════════════
# 12. SVG
# ═══════════════════════════════════════════════════════════════════════════════

def meta_svg(path):
    m = {}
    try:
        with open(path,"r",encoding="utf-8",errors="replace") as f:
            content = f.read()

        root = ET.fromstring(content)
        NS_SVG = "http://www.w3.org/2000/svg"
        NS_DC  = "http://purl.org/dc/elements/1.1/"
        NS_CC  = "http://creativecommons.org/ns#"
        NS_RDF = "http://www.w3.org/1999/02/22-rdf-syntax-ns#"

        # Atributos raíz
        m["svg_ancho"]    = root.get("width","")
        m["svg_alto"]     = root.get("height","")
        m["svg_viewbox"]  = root.get("viewBox","")
        m["svg_version"]  = root.get("version","")

        # Metadata embebida (Dublin Core / RDF)
        for el in root.iter():
            tag = el.tag.split("}")[-1] if "}" in el.tag else el.tag
            ns  = el.tag.split("}")[0].strip("{") if "}" in el.tag else ""
            if ns == NS_DC and el.text:
                m[f"svg_dc_{tag}"] = el.text.strip()
            if ns == NS_CC and el.text:
                m[f"svg_cc_{tag}"] = el.text.strip()

        m["svg_num_elementos"] = sum(1 for _ in root.iter())
        m["svg_longitud_chars"] = len(content)

    except Exception as e:
        m["svg_error"] = str(e)
    return m


# ═══════════════════════════════════════════════════════════════════════════════
# 13. EXIFTOOL (fallback universal + complemento)
# ═══════════════════════════════════════════════════════════════════════════════

def meta_exiftool(path):
    if not _exiftool_ok(): return {}
    try:
        r = subprocess.run(
            ["exiftool","-json","-a","-G1","-n","-long", str(path)],
            capture_output=True, text=True, check=True
        )
        data = json.loads(r.stdout)
        skip = {"SourceFile","ExifToolVersion","FileSize","FileModifyDate",
                "FileAccessDate","FileInodeChangeDate","FilePermissions",
                "FileType","FileTypeExtension","MIMEType","FileName","Directory"}
        result = {}
        for k, v in (data[0] if data else {}).items():
            clean_k = k.split(":")[-1] if ":" in k else k
            if clean_k in skip: continue
            result[f"et_{k.replace(':','_').replace(' ','_')}"] = _s(v)
        return result
    except: return {}


# ═══════════════════════════════════════════════════════════════════════════════
# DISPATCHER
# ═══════════════════════════════════════════════════════════════════════════════

DISPATCH = {
    # Imágenes
    **{e: meta_imagen for e in [
        ".jpg",".jpeg",".png",".tiff",".tif",".heic",".heif",
        ".webp",".bmp",".gif",".raw",".cr2",".nef",".arw",
        ".dng",".orf",".rw2",".pef",".srw",".x3f",".psd",
    ]},
    # PDF
    ".pdf": meta_pdf,
    # Office moderno
    ".docx": meta_docx, ".xlsx": meta_xlsx, ".pptx": meta_pptx,
    # Office antiguo
    ".doc": meta_doc_old, ".xls": meta_xls_old,
    # OpenDocument
    ".odt": lambda p: meta_opendoc(p,"odt"),
    ".ods": lambda p: meta_opendoc(p,"ods"),
    ".odp": lambda p: meta_opendoc(p,"odp"),
    # Audio
    **{e: meta_audio for e in [
        ".mp3",".flac",".ogg",".wav",".m4a",".aac",
        ".wma",".opus",".aiff",".aif",".ape",".mpc",".wv",
    ]},
    # Video
    **{e: meta_video for e in [
        ".mp4",".mov",".mkv",".avi",".webm",".flv",
        ".wmv",".m4v",".3gp",".ts",".mts",".m2ts",".vob",
    ]},
    # Otros
    ".epub": meta_epub,
    ".msg":  meta_msg,
    ".eml":  meta_eml,
    ".html": meta_html, ".htm": meta_html,
    ".zip":  meta_zip,
    ".svg":  meta_svg,
}


def extract_all(path):
    ext    = Path(path).suffix.lower()
    result = meta_filesystem(path)
    fn     = DISPATCH.get(ext)
    if fn:
        result.update(fn(path))
    # ExifTool SIEMPRE complementa (no sobreescribe, agrega lo que falta)
    et_data = meta_exiftool(path)
    for k, v in et_data.items():
        if k not in result:
            result[k] = v
    return result


# ═══════════════════════════════════════════════════════════════════════════════
# EXCEL
# ═══════════════════════════════════════════════════════════════════════════════

def build_excel(records, output_path):
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = openpyxl.Workbook()

    C_HDR_BG = "1F3864"; C_HDR_FG = "FFFFFF"
    C_ALT    = "EEF2F7"; C_BDR    = "B8C4D4"

    C_SEC = {
        "fs":    "D9E8FF",
        "img":   "FFF9C4", "exif": "FDE9D9", "iptc": "FCE4EC",
        "pdf":   "E8F5E9", "xmp":  "FFDDE8",
        "docx":  "F3E5F5", "xlsx": "E0F2F1", "pptx": "FFF3E0",
        "odt":   "E8EAF6", "ods":  "E8EAF6", "odp":  "E8EAF6",
        "xls":   "E0F2F1", "doc":  "F3E5F5",
        "audio": "E3F2FD",
        "video": "EDE7F6",
        "epub":  "E8F5E9",
        "msg":   "FBE9E7", "eml": "FBE9E7",
        "html":  "F9FBE7",
        "zip":   "ECEFF1",
        "svg":   "FFF8E1",
        "et":    "F5F5F5",
    }

    thin   = Side(style="thin", color=C_BDR)
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    def sec_color(k):
        for pfx, col in C_SEC.items():
            if k.startswith(pfx + "_"): return col
        return "FFFFFF"

    # Recopilar todas las columnas en orden de aparición
    all_keys = []
    seen = set()
    # Ordenar: primero fs_, luego por tipo, luego et_
    priority_prefixes = ["fs_","img_","exif_","iptc_","pdf_","xmp_",
                         "docx_","xlsx_","pptx_","odt_","ods_","odp_",
                         "doc_","xls_","audio_","video_","epub_",
                         "msg_","eml_","html_","zip_","svg_"]
    ordered = []
    buckets = {p: [] for p in priority_prefixes}
    rest = []
    for rec in records:
        for k in rec:
            if k not in seen:
                seen.add(k)
                placed = False
                for p in priority_prefixes:
                    if k.startswith(p):
                        buckets[p].append(k); placed = True; break
                if not placed:
                    rest.append(k)
    for p in priority_prefixes:
        ordered.extend(buckets[p])
    ordered.extend(rest)
    all_keys = ordered

    # ── Hoja 1: Metadatos completos ──────────────────────────────────────────
    ws = wb.active
    ws.title = "Metadatos"

    for ci, k in enumerate(all_keys, 1):
        c = ws.cell(row=1, column=ci, value=k)
        c.font      = Font(name="Arial", bold=True, color=C_HDR_FG, size=8)
        c.fill      = PatternFill("solid", fgColor=C_HDR_BG)
        c.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        c.border    = border
    ws.row_dimensions[1].height = 42
    ws.freeze_panes = "B2"

    for ri, rec in enumerate(records, 2):
        alt = (ri % 2 == 0)
        for ci, k in enumerate(all_keys, 1):
            val = rec.get(k, "")
            c   = ws.cell(row=ri, column=ci, value=val)
            c.font      = Font(name="Arial", size=8)
            c.border    = border
            c.alignment = Alignment(vertical="center", wrap_text=False)
            c.fill      = PatternFill("solid", fgColor=C_ALT if alt else sec_color(k))

    for ci, k in enumerate(all_keys, 1):
        max_w = max(len(k), 8)
        for rec in records:
            max_w = max(max_w, min(len(str(rec.get(k,""))), 55))
        ws.column_dimensions[get_column_letter(ci)].width = min(max_w + 2, 57)

    # ── Hoja 2: Resumen ──────────────────────────────────────────────────────
    ws2 = wb.create_sheet("Resumen")
    hdrs2 = [
        "Archivo","Extensión","Tamaño (KB)","Tipo MIME",
        "⚠ Fecha Modif. SO","Hash MD5",
        "✅ Fecha Creación Embebida","Autor / Creador",
        "Título","Dimensiones / Resolución",
        "Duración","Páginas / Hojas / Diapositivas",
        "GPS (decimal)","Link Google Maps",
        "Codec / Formato","Aplicación Creadora",
    ]
    for ci, h in enumerate(hdrs2, 1):
        c = ws2.cell(row=1, column=ci, value=h)
        c.font      = Font(name="Arial", bold=True, color=C_HDR_FG, size=9)
        c.fill      = PatternFill("solid", fgColor=C_HDR_BG)
        c.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        c.border    = border
    ws2.row_dimensions[1].height = 38
    ws2.freeze_panes = "A2"

    def _first(rec, *keys):
        for k in keys:
            v = rec.get(k,"")
            if v and str(v).strip(): return v
        return ""

    for ri, rec in enumerate(records, 2):
        alt  = (ri % 2 == 0)
        fill = PatternFill("solid", fgColor=C_ALT if alt else "FFFFFF")

        # Dimensiones
        dim = ""
        if rec.get("img_ancho_px"):
            dim = f"{rec['img_ancho_px']}x{rec['img_alto_px']} px"
            if rec.get("img_megapixeles"): dim += f" ({rec['img_megapixeles']} MP)"
        elif rec.get("video_v0_ancho_px"):
            dim = f"{rec['video_v0_ancho_px']}x{rec['video_v0_alto_px']} px"
        elif rec.get("img_dpi"):
            dim = f"DPI: {rec['img_dpi']}"

        vals = [
            rec.get("fs_nombre_archivo",""),
            rec.get("fs_extension",""),
            rec.get("fs_tamaño_kb",""),
            rec.get("fs_tipo_mime",""),
            rec.get("fs_fecha_modificacion",""),
            rec.get("fs_md5",""),
            _first(rec,
                "exif_DateTimeOriginal","exif_Exif_DateTimeOriginal",
                "exif_0th_DateTime","exif_Exif_DateTimeDigitized",
                "pdf_fecha_creacion","docx_fecha_creacion","docx_created",
                "xlsx_fecha_creacion","xlsx_created",
                "pptx_fecha_creacion","pptx_created",
                "odt_fecha_creacion","ods_fecha_creacion","odp_fecha_creacion",
                "audio_date","audio_tag_tdrc","audio_tag_year",
                "video_tag_creation_time","video_tag_date",
                "epub_fecha","msg_fecha_envio","eml_fecha_envio",
                "et_EXIF_DateTimeOriginal","et_QuickTime_CreateDate",
                "et_ID3_Year","et_XMP_CreateDate",
            ),
            _first(rec,
                "exif_Artist","iptc_autor",
                "docx_autor_creador","docx_author","xlsx_creator","xlsx_autor_creador",
                "pptx_author","pptx_autor_creador",
                "odt_creador_original","ods_creador_original","odp_creador_original",
                "pdf_author","audio_artist","audio_tag_tpe1",
                "video_tag_artist","epub_autor",
                "msg_remitente","eml_remitente",
                "et_EXIF_Artist","et_ID3_Artist","et_XMP_Creator",
            ),
            _first(rec,
                "docx_titulo","xlsx_titulo","pptx_titulo",
                "pdf_title","audio_title","audio_tag_tit2",
                "video_tag_title","epub_titulo",
                "msg_asunto","eml_asunto","html_titulo",
                "exif_ImageDescription","et_XMP_Title","et_ID3_Title",
            ),
            dim,
            _first(rec,"audio_duracion_hms","video_duracion_hms",
                       "audio_duracion_seg","video_duracion_seg"),
            _first(rec,"pdf_num_paginas","xlsx_num_hojas",
                       "pptx_num_diapositivas","epub_num_capitulos",
                       "docx_app_pages","docx_app_slides"),
            _first(rec,"exif_GPS_latitud_decimal","et_GPS_Latitude"),
            _first(rec,"exif_GPS_google_maps"),
            _first(rec,"audio_codec","video_formato_corto",
                       "video_v0_codec","img_formato","et_FileType"),
            _first(rec,
                "docx_app_application","xlsx_app_application","pptx_app_application",
                "odt_aplicacion_creadora","ods_aplicacion_creadora",
                "pdf_creator","pdf_producer",
                "exif_Software","et_XMP_CreatorTool","et_Software",
            ),
        ]
        for ci, val in enumerate(vals, 1):
            c = ws2.cell(row=ri, column=ci, value=val)
            c.font      = Font(name="Arial", size=8)
            c.border    = border
            c.alignment = Alignment(vertical="center")
            c.fill      = fill

    ws2_widths = [30,10,11,24,22,34,26,24,28,22,14,12,18,32,16,24]
    for ci, w in enumerate(ws2_widths, 1):
        ws2.column_dimensions[get_column_letter(ci)].width = w

    # ── Hoja 3: Leyenda ──────────────────────────────────────────────────────
    ws3 = wb.create_sheet("Leyenda")
    ws3["A1"] = "Leyenda de prefijos de columnas"
    ws3["A1"].font = Font(name="Arial", bold=True, size=12)

    leyenda = [
        ("fs_",    "Sistema de archivos — puede cambiar al descargar (NO embebido)",   C_SEC["fs"]),
        ("img_",   "Propiedades de imagen: formato, dimensiones, DPI, ICC, GIF frames", C_SEC["img"]),
        ("exif_",  "Datos EXIF embebidos: fecha original, cámara, GPS, exposición, lente", C_SEC["exif"]),
        ("iptc_",  "Datos IPTC: autor, copyright, palabras clave, localización",        C_SEC["iptc"]),
        ("xmp_",   "Bloque XMP (Adobe/estándar abierto): derechos, historial, ID",      C_SEC["xmp"]),
        ("pdf_",   "Metadatos internos PDF: autor, fecha, páginas, fuentes, permisos",  C_SEC["pdf"]),
        ("docx_",  "Propiedades Word: autor, fecha creación, palabras, revisiones",     C_SEC["docx"]),
        ("xlsx_",  "Propiedades Excel: autor, fecha, hojas, estadísticas",              C_SEC["xlsx"]),
        ("pptx_",  "Propiedades PowerPoint: autor, fecha, diapositivas, layouts",       C_SEC["pptx"]),
        ("odt_",   "OpenDocument Texto: autor, fecha, estadísticas",                   C_SEC["odt"]),
        ("ods_",   "OpenDocument Hoja de cálculo: autor, fecha",                        C_SEC["ods"]),
        ("odp_",   "OpenDocument Presentación: autor, fecha",                           C_SEC["odp"]),
        ("doc_",   "Word binario antiguo (DOC): formato, metadata básica",              C_SEC["doc"]),
        ("xls_",   "Excel binario antiguo (XLS): hojas, autor",                        C_SEC["xls"]),
        ("audio_", "Tags embebidos de audio: ID3/Vorbis — artista, álbum, año, ISRC",  C_SEC["audio"]),
        ("video_", "Streams y tags embebidos de video (ffprobe): codec, GPS, rotación", C_SEC["video"]),
        ("epub_",  "Metadatos EPUB: título, autor, editorial, fecha, idioma",           C_SEC["epub"]),
        ("msg_",   "Propiedades Outlook MSG: asunto, remitente, fecha, adjuntos",       C_SEC["msg"]),
        ("eml_",   "Cabeceras EML: asunto, remitente, IP origen, ruta servidores",     C_SEC["eml"]),
        ("html_",  "Meta tags HTML: título, description, og:*, charset",               C_SEC["html"]),
        ("zip_",   "Contenido ZIP: archivos, fechas embebidas, ratio compresión",       C_SEC["zip"]),
        ("svg_",   "Metadatos SVG: dimensiones, Dublin Core embebido",                  C_SEC["svg"]),
        ("et_",    "ExifTool: campos adicionales (MakerNote, propietarios, etc.)",      C_SEC["et"]),
    ]

    ws3.row_dimensions[1].height = 22
    for i, (pfx, desc, color) in enumerate(leyenda, 3):
        c1 = ws3.cell(row=i, column=1, value=pfx)
        c2 = ws3.cell(row=i, column=2, value=desc)
        for c in (c1, c2):
            c.fill      = PatternFill("solid", fgColor=color)
            c.font      = Font(name="Arial", size=9)
            c.border    = border
            c.alignment = Alignment(vertical="center", wrap_text=False)
        ws3.row_dimensions[i].height = 16

    ws3.column_dimensions["A"].width = 12
    ws3.column_dimensions["B"].width = 70

    ws3.cell(row=len(leyenda)+5, column=1,
             value="⚠ IMPORTANTE: Los campos fs_ dependen del sistema operativo y CAMBIAN al copiar o descargar el archivo.").font = Font(
        name="Arial", bold=True, color="CC0000", size=9)

    wb.save(output_path)
    print(f"\n✅ Excel guardado en: {output_path}")


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def scan_folder(folder, recursive=False):
    p = Path(folder)
    pattern = "**/*" if recursive else "*"
    return sorted(str(f) for f in p.glob(pattern) if f.is_file())

def main():
    parser = argparse.ArgumentParser(
        description="Extrae TODOS los metadatos posibles de una carpeta → Excel",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Ejemplos:
  python meta.py C:/fotos    C:/resultados
  python meta.py ./archivos  ./salida  --recursivo
  python meta.py ./archivos  ./salida  --recursivo --nombre inventario_mayo
        """
    )
    parser.add_argument("entrada",  help="Carpeta con los archivos a analizar")
    parser.add_argument("salida",   help="Carpeta donde se guardará el Excel")
    parser.add_argument("--recursivo","-r", action="store_true",
                        help="Incluir subcarpetas")
    parser.add_argument("--nombre","-n", default="",
                        help="Nombre del Excel de salida (sin .xlsx)")
    args = parser.parse_args()

    entrada = Path(args.entrada)
    salida  = Path(args.salida)

    if not entrada.is_dir():
        print(f"[ERROR] No existe: {entrada}"); sys.exit(1)

    salida.mkdir(parents=True, exist_ok=True)
    ts     = datetime.now().strftime("%Y%m%d_%H%M%S")
    nombre = args.nombre or f"metadatos_{ts}"
    xlsx   = salida / f"{nombre}.xlsx"

    archivos = scan_folder(entrada, args.recursivo)
    if not archivos:
        print("No se encontraron archivos."); sys.exit(0)

    print(f"\n📂 {len(archivos)} archivos encontrados en '{entrada}'")
    if _exiftool_ok(): print("   ✅ ExifTool detectado — extracción máxima activada")
    else:              print("   ⚠  ExifTool NO detectado — instálalo para más datos")
    if _ffprobe_ok():  print("   ✅ FFmpeg/ffprobe detectado — videos soportados")
    else:              print("   ⚠  FFmpeg NO detectado — los videos no se analizarán")
    print()

    records = []
    for i, f in enumerate(archivos, 1):
        print(f"  [{i:>4}/{len(archivos)}] {Path(f).name}")
        try:    records.append(extract_all(f))
        except Exception as e:
            print(f"          ⚠ Error: {e}")
            records.append({"fs_nombre_archivo": Path(f).name, "error_general": str(e)})

    print(f"\n📊 Generando Excel ({len(records)} archivos, {len(set(k for r in records for k in r))} columnas)...")
    build_excel(records, str(xlsx))

if __name__ == "__main__":
    main()
